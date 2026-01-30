from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import json
import os
from pathlib import Path

class BrandConfig:
    """Loads and manages Zain brand configuration"""
    
    def __init__(self, config_path="config.json"):
        with open(config_path, 'r') as f:
            self.config = json.load(f)
        self.base_path = Path(config_path).parent
    
    def get_color(self, name):
        """Get hex color by name"""
        colors = self.config.get("colors", {})
        if name in colors:
            return colors[name]
        # Check gradients
        gradients = colors.get("gradients", {})
        if name in gradients:
            return gradients[name]
        return None
    
    def hex_to_rgb(self, hex_color):
        """Convert hex to RGB tuple"""
        hex_color = hex_color.lstrip('#')
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
    
    def get_logo_path(self, variant="english_black"):
        """Get absolute path to logo"""
        logo_rel = self.config["logos"].get(variant)
        if logo_rel:
            return self.base_path / logo_rel
        return None
    
    def get_font_name(self, style="primary"):
        """Get font name for style"""
        return self.config["fonts"].get(style, "Zain-Regular")
    
    def get_slide_layout(self, layout_type):
        """Get layout configuration"""
        return self.config["slide_layouts"].get(layout_type, {})
    
    def get_gradient_path(self, gradient_name):
        """Get path to gradient background image"""
        gradient_map = {
            "ultraviolet": "brand-assets/gradients/ZN_GRD_16x9_ULTRAVIOLET.png",
            "limelagoon": "brand-assets/gradients/ZN_GRD_16x9_LIMELAGOON.png",
            "midnightsky": "brand-assets/gradients/ZN_GRD_16x9_MIDNIGHTSKY.png",
            "coraldawn": "brand-assets/gradients/ZN_GRD_16x9_CORALDAWN.png",
            "twilightmist": "brand-assets/gradients/ZN_GRD_16x9_TWILIGHTMIST.png",
            "azurewaters": "brand-assets/gradients/ZN_GRD_16x9_AZUREWATERS.png",
            "magentafade": "brand-assets/gradients/ZN_GRD_16x9_MAGENTAFADE.png",
            "jadehorizon": "brand-assets/gradients/ZN_GRD_16x9_JADEHORIZON.png"
        }
        gradient_rel = gradient_map.get(gradient_name)
        if gradient_rel:
            return self.base_path / gradient_rel
        return None


class SlideBuilder:
    """Builds individual slides with Zain branding"""
    
    def __init__(self, brand_config):
        self.brand = brand_config
    
    def add_title_slide(self, prs, title, subtitle="", gradient="ultraviolet"):
        """Create title slide with gradient background"""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add gradient background image
        gradient_path = self.brand.get_gradient_path(gradient)
        if gradient_path and gradient_path.exists():
            slide.shapes.add_picture(
                str(gradient_path), Inches(0), Inches(0),
                width=Inches(13.333), height=Inches(7.5)
            )
        else:
            # Fallback to solid color
            background = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), 
                Inches(13.333), Inches(7.5)
            )
            background.fill.solid()
            background.fill.fore_color.rgb = RGBColor(0x6E, 0x2C, 0x91)
            background.line.fill.background()
        
        # Add logo
        logo_path = self.brand.get_logo_path("english_white")
        if logo_path and logo_path.exists():
            slide.shapes.add_picture(
                str(logo_path), Inches(10), Inches(0.5), 
                width=Inches(2.5)
            )
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(2.5), Inches(11.8), Inches(1.5)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)  # White
        p.font.name = self.brand.get_font_name("bold")
        
        # Add subtitle
        if subtitle:
            sub_box = slide.shapes.add_textbox(
                Inches(0.75), Inches(4.2), Inches(11.8), Inches(1)
            )
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.font.name = self.brand.get_font_name("primary")
        
        return slide
    
    def add_content_slide(self, prs, title, bullets):
        """Create content slide with bullet points"""
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)
        
        # White background
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
            Inches(13.333), Inches(7.5)
        )
        background.fill.solid()
        background.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        background.line.fill.background()
        
        # Add logo (black version for white background)
        logo_path = self.brand.get_logo_path("english_black")
        if logo_path and logo_path.exists():
            slide.shapes.add_picture(
                str(logo_path), Inches(10.5), Inches(6.5),
                width=Inches(2)
            )
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(0.5), Inches(11.8), Inches(1)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x6E, 0x2C, 0x91)  # Primary purple
        p.font.name = self.brand.get_font_name("bold")
        
        # Add bullet points
        content_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(1.8), Inches(11.8), Inches(4.5)
        )
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {bullet}"
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
            p.font.name = self.brand.get_font_name("primary")
            p.space_after = Pt(12)
        
        return slide
    
    def add_section_slide(self, prs, title, gradient="coraldawn"):
        """Create section divider slide"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # Add gradient background image
        gradient_path = self.brand.get_gradient_path(gradient)
        if gradient_path and gradient_path.exists():
            slide.shapes.add_picture(
                str(gradient_path), Inches(0), Inches(0),
                width=Inches(13.333), height=Inches(7.5)
            )
        else:
            # Fallback to solid color
            background = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                Inches(13.333), Inches(7.5)
            )
            background.fill.solid()
            background.fill.fore_color.rgb = RGBColor(0xE6, 0x00, 0x7E)
            background.line.fill.background()
        
        # Add white logo
        logo_path = self.brand.get_logo_path("english_white")
        if logo_path and logo_path.exists():
            slide.shapes.add_picture(
                str(logo_path), Inches(10), Inches(0.5),
                width=Inches(2.5)
            )
        
        # Large section title
        title_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(2.5), Inches(11.8), Inches(2)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.font.name = self.brand.get_font_name("black")
        
        return slide


class PPTXAssembler:
    """Assembles final presentation from generated slides"""
    
    def __init__(self, brand_config):
        self.brand = brand_config
        self.builder = SlideBuilder(brand_config)
    
    def create_presentation(self, slides_data, output_path):
        """
        Create full presentation from slide data
        
        slides_data: list of dicts with 'type', 'title', 'content'
        """
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        for slide_data in slides_data:
            slide_type = slide_data.get('type', 'content')
            title = slide_data.get('title', '')
            content = slide_data.get('content', [])
            
            if slide_type == 'title':
                subtitle = slide_data.get('subtitle', '')
                gradient = slide_data.get('gradient', 'ultraviolet')
                self.builder.add_title_slide(prs, title, subtitle, gradient)
            
            elif slide_type == 'section':
                gradient = slide_data.get('gradient', 'coraldawn')
                self.builder.add_section_slide(prs, title, gradient)
            
            elif slide_type == 'content':
                if isinstance(content, str):
                    content = [content]
                self.builder.add_content_slide(prs, title, content)
        
        prs.save(output_path)
        return output_path
