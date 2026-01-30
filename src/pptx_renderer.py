"""
AI-Powered Slide Generator
Uses sub-agents to analyze content and generate intelligent presentations
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path
import json


class BrandConfig:
    """Loads and manages Zain brand configuration"""
    
    def __init__(self, config_path="config.json"):
        with open(config_path, 'r') as f:
            self.config = json.load(f)
        self.base_path = Path(config_path).parent
    
    def get_logo_path(self, variant="english_black"):
        logo_rel = self.config["logos"].get(variant)
        if logo_rel:
            return self.base_path / logo_rel
        return None
    
    def get_font_name(self, style="primary"):
        return self.config["fonts"].get(style, "Zain-Regular")
    
    def get_gradient_path(self, gradient_name):
        gradient_map = {
            "ultraviolet": "brand-assets/gradients/ZN_GRD_16x9_ULTRAVIOLET.png",
            "limelagoon": "brand-assets/gradients/ZN_GRD_16x9_LIMELAGOON.png",
            "midnightsky": "brand-assets/gradients/ZN_GRD_16x9_MIDNIGHTSKY.png",
            "coraldawn": "brand-assets/gradients/ZN_GRD_16x9_CORALDAWN.png",
            "twilightmist": "brand-assets/gradients/ZN_GRD_16x9_TWILIGHTMIST.png",
            "azurewaters": "brand-assets/gradients/ZN_GRD_16x9_AZUREWATERS.png",
            "magentafade": "brand-assets/gradients/ZN_GRD_16x9_MAGENTAFADE.png",
            "jadehorizon": "brand-assets/gradients/ZN_GRD_16x9_JADEHORIZON.png"
        }
        gradient_rel = gradient_map.get(gradient_name)
        if gradient_rel:
            return self.base_path / gradient_rel
        return None


class SlideRenderer:
    """Renders slide data to PPTX (no AI, just rendering)"""
    
    def __init__(self, brand_config):
        self.brand = brand_config
    
    def render_title_slide(self, slide, title, subtitle, gradient="ultraviolet"):
        """Render a title slide"""
        # Add gradient background
        gradient_path = self.brand.get_gradient_path(gradient)
        if gradient_path and gradient_path.exists():
            slide.shapes.add_picture(
                str(gradient_path), Inches(0), Inches(0),
                width=Inches(13.333), height=Inches(7.5)
            )
        
        # Add logo
        logo_path = self.brand.get_logo_path("english_white")
        if logo_path and logo_path.exists():
            slide.shapes.add_picture(
                str(logo_path), Inches(10), Inches(0.5), 
                width=Inches(2.5)
            )
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(2.5), Inches(11.8), Inches(1.5)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.font.name = self.brand.get_font_name("bold")
        
        # Add subtitle
        if subtitle:
            sub_box = slide.shapes.add_textbox(
                Inches(0.75), Inches(4.2), Inches(11.8), Inches(1)
            )
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.font.name = self.brand.get_font_name("primary")
    
    def render_content_slide(self, slide, title, content_items, layout="bullets"):
        """Render a content slide with various layouts"""
        # White background
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
            Inches(13.333), Inches(7.5)
        )
        background.fill.solid()
        background.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        background.line.fill.background()
        
        # Add logo
        logo_path = self.brand.get_logo_path("english_black")
        if logo_path and logo_path.exists():
            slide.shapes.add_picture(
                str(logo_path), Inches(10.5), Inches(6.5),
                width=Inches(2)
            )
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(0.5), Inches(11.8), Inches(1)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x6E, 0x2C, 0x91)
        p.font.name = self.brand.get_font_name("bold")
        
        # Render content based on layout
        if layout == "bullets":
            self._render_bullets(slide, content_items)
        elif layout == "two_column":
            self._render_two_column(slide, content_items)
        elif layout == "big_number":
            self._render_big_number(slide, content_items)
        elif layout == "quote":
            self._render_quote(slide, content_items)
    
    def _render_bullets(self, slide, items):
        """Render bullet points"""
        content_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(1.8), Inches(11.8), Inches(4.5)
        )
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
            p.font.name = self.brand.get_font_name("primary")
            p.space_after = Pt(12)
    
    def _render_two_column(self, slide, items):
        """Render two-column layout"""
        mid = len(items) // 2
        left_items = items[:mid] if mid > 0 else items
        right_items = items[mid:] if mid > 0 else []
        
        # Left column
        left_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(1.8), Inches(5.5), Inches(4.5)
        )
        tf = left_box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(left_items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
            p.font.name = self.brand.get_font_name("primary")
            p.space_after = Pt(10)
        
        # Right column
        if right_items:
            right_box = slide.shapes.add_textbox(
                Inches(7), Inches(1.8), Inches(5.5), Inches(4.5)
            )
            tf = right_box.text_frame
            tf.word_wrap = True
            for i, item in enumerate(right_items):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(16)
                p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
                p.font.name = self.brand.get_font_name("primary")
                p.space_after = Pt(10)
    
    def _render_big_number(self, slide, items):
        """Render slide with big number/stat highlight"""
        if items:
            # First item is the big number
            big_num_box = slide.shapes.add_textbox(
                Inches(0.75), Inches(2), Inches(6), Inches(2)
            )
            tf = big_num_box.text_frame
            p = tf.paragraphs[0]
            p.text = items[0]
            p.font.size = Pt(60)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0x6E, 0x2C, 0x91)
            p.font.name = self.brand.get_font_name("black")
            
            # Rest are supporting points
            if len(items) > 1:
                support_box = slide.shapes.add_textbox(
                    Inches(7), Inches(2), Inches(5.5), Inches(4)
                )
                tf = support_box.text_frame
                tf.word_wrap = True
                for i, item in enumerate(items[1:]):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = f"• {item}"
                    p.font.size = Pt(16)
                    p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
                    p.font.name = self.brand.get_font_name("primary")
                    p.space_after = Pt(8)
    
    def _render_quote(self, slide, items):
        """Render a quote-style slide"""
        if items:
            quote_box = slide.shapes.add_textbox(
                Inches(1), Inches(2), Inches(11.3), Inches(3.5)
            )
            tf = quote_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f'"{items[0]}"'
            p.font.size = Pt(28)
            p.font.italic = True
            p.font.color.rgb = RGBColor(0x6E, 0x2C, 0x91)
            p.font.name = self.brand.get_font_name("italic")
            
            if len(items) > 1:
                attr_box = slide.shapes.add_textbox(
                    Inches(1), Inches(5.5), Inches(11.3), Inches(0.8)
                )
                tf = attr_box.text_frame
                p = tf.paragraphs[0]
                p.text = f"— {items[1]}"
                p.font.size = Pt(18)
                p.font.color.rgb = RGBColor(0x4A, 0x4A, 0x4A)
                p.font.name = self.brand.get_font_name("primary")
    
    def render_section_slide(self, slide, title, gradient="coraldawn"):
        """Render a section divider slide"""
        gradient_path = self.brand.get_gradient_path(gradient)
        if gradient_path and gradient_path.exists():
            slide.shapes.add_picture(
                str(gradient_path), Inches(0), Inches(0),
                width=Inches(13.333), height=Inches(7.5)
            )
        
        logo_path = self.brand.get_logo_path("english_white")
        if logo_path and logo_path.exists():
            slide.shapes.add_picture(
                str(logo_path), Inches(10), Inches(0.5),
                width=Inches(2.5)
            )
        
        title_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(2.5), Inches(11.8), Inches(2)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.font.name = self.brand.get_font_name("black")


class PPTXRenderer:
    """Main renderer that takes AI-generated slide data and creates PPTX"""
    
    def __init__(self, brand_config):
        self.brand = brand_config
        self.slide_renderer = SlideRenderer(brand_config)
    
    def render_presentation(self, ai_slides_data, output_path):
        """
        Render AI-generated slides to PPTX
        
        ai_slides_data: List of slide dicts with:
            - type: 'title' | 'section' | 'content'
            - title: slide title
            - subtitle: (for title slides)
            - content: list of content items
            - layout: 'bullets' | 'two_column' | 'big_number' | 'quote'
            - gradient: gradient name
        """
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        for slide_data in ai_slides_data:
            slide_type = slide_data.get('type', 'content')
            slide_layout = prs.slide_layouts[6]  # Blank
            slide = prs.slides.add_slide(slide_layout)
            
            if slide_type == 'title':
                self.slide_renderer.render_title_slide(
                    slide,
                    slide_data.get('title', ''),
                    slide_data.get('subtitle', ''),
                    slide_data.get('gradient', 'ultraviolet')
                )
            
            elif slide_type == 'section':
                self.slide_renderer.render_section_slide(
                    slide,
                    slide_data.get('title', ''),
                    slide_data.get('gradient', 'coraldawn')
                )
            
            elif slide_type == 'content':
                self.slide_renderer.render_content_slide(
                    slide,
                    slide_data.get('title', ''),
                    slide_data.get('content', []),
                    slide_data.get('layout', 'bullets')
                )
        
        prs.save(output_path)
        return output_path
